"""
Enhanced PPTX Renderer Service V2

Integrates intelligent layout selection with visual rendering.
Provides professional slide generation with automatic layout optimization.

Created: 2026-03-23
"""

from typing import List, Dict, Any, Optional
from pathlib import Path
from loguru import logger
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Import layout modules
from lib.layout import (
    get_layout_selector,
    get_template_manager,
    ColorScheme,
    Position,
    LayoutRenderer,
    IntelligentLayoutSelector,
    TemplateManager,
)


class PPTXRendererV2:
    """
    Enhanced PPTX Renderer with intelligent layout selection.

    Features:
    - Automatic layout selection based on content analysis
    - Theme integration (5 professional themes)
    - Visual element rendering (charts, shapes, diagrams)
    - Support for 40+ layout configurations
    """

    # Slide dimensions (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Content area (excluding margins)
    CONTENT_MARGIN_LEFT = 0.5
    CONTENT_MARGIN_RIGHT = 0.5
    CONTENT_MARGIN_TOP = 1.5
    CONTENT_MARGIN_BOTTOM = 0.8

    def __init__(
        self,
        theme_id: str = "blue_professional",
        output_dir: str = "./output/pptx",
        template_path: Optional[str] = None,
        auto_layout: bool = True
    ):
        """
        Initialize the enhanced renderer.

        Args:
            theme_id: Theme identifier (blue_professional, green_natural, etc.)
            output_dir: Output directory for generated files
            template_path: Optional PPTX template file
            auto_layout: Enable automatic layout selection
        """
        self.theme_id = theme_id
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.template_path = template_path
        self.auto_layout = auto_layout

        # Initialize components
        self.prs = None
        self.template_manager = get_template_manager()
        self.layout_selector = get_layout_selector()
        self.layout_renderer = None  # Initialized per theme

        # Load theme
        self._load_theme(theme_id)

        logger.info(f"PPTXRendererV2 initialized (theme: {theme_id}, auto_layout: {auto_layout})")

    def _load_theme(self, theme_id: str):
        """Load theme configuration"""
        theme = self.template_manager.get_theme(theme_id)

        if theme:
            self.colors = ColorScheme(
                primary=RGBColor(
                    int(theme.primary_color[1:3], 16),
                    int(theme.primary_color[3:5], 16),
                    int(theme.primary_color[5:7], 16)
                ),
                secondary=RGBColor(
                    int(theme.secondary_color[1:3], 16),
                    int(theme.secondary_color[3:5], 16),
                    int(theme.secondary_color[5:7], 16)
                ),
                accent=RGBColor(
                    int(theme.accent_color[1:3], 16),
                    int(theme.accent_color[3:5], 16),
                    int(theme.accent_color[5:7], 16)
                ),
                text_dark=RGBColor(51, 51, 51),
                text_light=RGBColor(128, 128, 128),
                background=RGBColor(245, 245, 245)
            )
            self.font_family = theme.font_family
        else:
            # Fallback to default blue theme
            self.colors = ColorScheme(
                primary=RGBColor(0, 82, 139),
                secondary=RGBColor(51, 102, 153),
                accent=RGBColor(231, 76, 60),
                text_dark=RGBColor(51, 51, 51),
                text_light=RGBColor(128, 128, 128),
                background=RGBColor(245, 245, 245)
            )
            self.font_family = "Microsoft YaHei"

        # Initialize layout renderer with colors
        self.layout_renderer = LayoutRenderer(self.colors)

    def render_report(
        self,
        slides: List[Dict[str, Any]],
        report_id: str,
        client_name: str,
        use_intelligent_selection: bool = True
    ) -> str:
        """
        Render complete presentation with intelligent layout selection.

        Args:
            slides: List of slide data dicts
            report_id: Report identifier
            client_name: Client name
            use_intelligent_selection: Use AI-powered layout selection

        Returns:
            Path to generated PPTX file
        """
        output_path = self.output_dir / f"{report_id}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"

        # Create presentation
        if self.template_path and Path(self.template_path).exists():
            self.prs = Presentation(self.template_path)
        else:
            self.prs = Presentation()

        # Set slide dimensions
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT

        logger.info(f"Creating presentation: {output_path}")
        logger.info(f"Total slides: {len(slides)}, intelligent_selection: {use_intelligent_selection}")

        try:
            for i, slide_data in enumerate(slides):
                # Select layout intelligently if enabled
                if use_intelligent_selection and self.auto_layout:
                    slide_data = self._enhance_slide_with_layout(slide_data)

                self._render_slide(slide_data, i + 1)

            # Save
            self.prs.save(output_path)
            logger.info(f"Report saved: {output_path}")

            return str(output_path)

        except Exception as e:
            logger.error(f"Error rendering report: {e}")
            raise

    def _enhance_slide_with_layout(self, slide_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Enhance slide data with intelligent layout selection.

        Uses ContentAnalyzer + IntelligentLayoutSelector to pick optimal layout.
        """
        # Skip if layout already specified
        if slide_data.get("layout") and slide_data.get("layout") != "auto":
            return slide_data

        try:
            # Use intelligent selector
            result = self.layout_selector.select(
                title=slide_data.get("title", ""),
                key_message=slide_data.get("key_message", ""),
                bullets=slide_data.get("bullets", []),
                context=slide_data.get("context", {})
            )

            # Add layout info to slide data
            slide_data["layout"] = result.primary.layout_id
            slide_data["layout_confidence"] = result.primary.confidence
            slide_data["layout_reason"] = result.primary.reason
            slide_data["layout_params"] = result.primary.params
            slide_data["content_analysis"] = result.analysis.to_dict()

            logger.debug(
                f"Selected layout {result.primary.layout_id} "
                f"(confidence: {result.primary.confidence:.2f}) "
                f"for slide: {slide_data.get('title', '')[:30]}"
            )

        except Exception as e:
            logger.warning(f"Layout selection failed, using default: {e}")
            slide_data["layout"] = "BULLET_01"

        return slide_data

    def _render_slide(self, slide_data: Dict[str, Any], slide_index: int):
        """
        Render a single slide.

        Args:
            slide_data: Slide content and layout info
            slide_index: 1-based slide index
        """
        # Add blank slide
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        layout_id = slide_data.get("layout", "bullet_points")

        # Common elements: title and key_message
        self._render_slide_header(slide, slide_data)

        # Route to layout-specific renderer
        if layout_id in ["title_slide", "TITLE_01"]:
            self._render_title_slide(slide, slide_data)
        elif layout_id in ["section_divider", "SECTION_01"]:
            self._render_section_divider(slide, slide_data)
        elif layout_id.startswith("KEY_INSIGHT"):
            self._render_key_insight(slide, slide_data)
        else:
            # Use LayoutRenderer for content slides
            self._render_content_slide(slide, slide_data)

        # Add source reference if present
        self._render_source_ref(slide, slide_data)

    def _render_slide_header(self, slide, slide_data: Dict[str, Any]):
        """Render slide title and key message"""
        # Title
        title = slide_data.get("title", "")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.colors.text_dark

        # Key message (action title)
        key_message = slide_data.get("key_message", "")
        if key_message:
            msg_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.5)
            )
            tf = msg_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = key_message
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors.primary

    def _render_title_slide(self, slide, slide_data: Dict[str, Any]):
        """Render title/cover slide"""
        # Clear any header elements added
        for shape in list(slide.shapes):
            if shape.top < Inches(2.0):
                sp = shape._element
                sp.getparent().remove(sp)

        # Main title - centered
        title = slide_data.get("title", "咨询项目建议书")
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.colors.primary
        p.alignment = PP_ALIGN.CENTER

        # Subtitle / client name
        subtitle = slide_data.get("subtitle", slide_data.get("client_name", ""))
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = self.colors.text_dark
            p.alignment = PP_ALIGN.CENTER

        # Date
        date_str = slide_data.get("date", datetime.now().strftime("%Y年%m月"))
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = date_str
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors.text_light
        p.alignment = PP_ALIGN.CENTER

        # Decorative line
        line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(4.5), Inches(4.0), Inches(4.333), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors.accent
        line.line.fill.background()

    def _render_section_divider(self, slide, slide_data: Dict[str, Any]):
        """Render section divider slide"""
        # Clear header elements
        for shape in list(slide.shapes):
            if shape.top < Inches(2.0):
                sp = shape._element
                sp.getparent().remove(sp)

        # Section number/badge
        section_num = slide_data.get("section_num", "")
        if section_num:
            from lib.layout.visual_elements import ShapeFactory, Position
            badge_pos = Position(0.5, 2.0, 0.8, 0.8)
            ShapeFactory.add_circle(
                slide, badge_pos, self.colors.accent, section_num, font_size=24
            )

        # Section title
        title = slide_data.get("title", slide_data.get("subsection", ""))
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.0), Inches(11.333), Inches(1.2)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = self.colors.primary

        # Description
        desc = slide_data.get("description", slide_data.get("key_message", ""))
        if desc:
            desc_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(3.5), Inches(11.333), Inches(1.0)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors.text_light

    def _render_key_insight(self, slide, slide_data: Dict[str, Any]):
        """Render key insight slide - large centered message"""
        # Use LayoutRenderer
        position = Position(
            self.CONTENT_MARGIN_LEFT,
            self.CONTENT_MARGIN_TOP,
            self.SLIDE_WIDTH - Inches(self.CONTENT_MARGIN_LEFT + self.CONTENT_MARGIN_RIGHT),
            self.SLIDE_HEIGHT - Inches(self.CONTENT_MARGIN_TOP + self.CONTENT_MARGIN_BOTTOM)
        )

        self.layout_renderer.render(slide, "KEY_INSIGHT_01", slide_data, position)

    def _render_content_slide(self, slide, slide_data: Dict[str, Any]):
        """Render content slide using LayoutRenderer"""
        layout_id = slide_data.get("layout", "BULLET_01")

        # Calculate content area
        position = Position(
            self.CONTENT_MARGIN_LEFT,
            self.CONTENT_MARGIN_TOP,
            self.SLIDE_WIDTH - Inches(self.CONTENT_MARGIN_LEFT + self.CONTENT_MARGIN_RIGHT),
            self.SLIDE_HEIGHT - Inches(self.CONTENT_MARGIN_TOP + self.CONTENT_MARGIN_BOTTOM)
        )

        # Apply layout params if present
        layout_params = slide_data.get("layout_params", {})
        if layout_params:
            # Adjust font sizes based on text density
            pass

        # Render using LayoutRenderer
        self.layout_renderer.render(slide, layout_id, slide_data, position)

    def _render_source_ref(self, slide, slide_data: Dict[str, Any]):
        """Render source reference at bottom"""
        source_ref = slide_data.get("source_ref", "")
        if source_ref:
            source_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4)
            )
            tf = source_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"来源: {source_ref}"
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors.text_light

    def set_theme(self, theme_id: str):
        """Change theme mid-presentation"""
        self.theme_id = theme_id
        self._load_theme(theme_id)
        logger.info(f"Theme changed to: {theme_id}")


# Convenience function
def create_presentation_v2(
    slides: List[Dict[str, Any]],
    report_id: str,
    client_name: str,
    output_dir: str = "./output/pptx",
    theme_id: str = "blue_professional",
    auto_layout: bool = True
) -> str:
    """
    Create presentation with V2 renderer.

    Args:
        slides: Slide data list
        report_id: Report ID
        client_name: Client name
        output_dir: Output directory
        theme_id: Theme identifier
        auto_layout: Enable automatic layout selection

    Returns:
        Path to generated PPTX
    """
    renderer = PPTXRendererV2(
        theme_id=theme_id,
        output_dir=output_dir,
        auto_layout=auto_layout
    )
    return renderer.render_report(slides, report_id, client_name)


# Integration with workflow
def render_slides_from_drafts(
    slide_drafts: List[Any],
    report_id: str,
    client_name: str,
    theme_id: str = "blue_professional"
) -> str:
    """
    Render slides from SlideDraft objects.

    This function integrates with the LangGraph workflow output.

    Args:
        slide_drafts: List of SlideDraft objects (from workflow)
        report_id: Report ID
        client_name: Client name
        theme_id: Theme to use

    Returns:
        Path to generated PPTX
    """
    # Convert SlideDraft to dict
    slides = []
    for draft in slide_drafts:
        if hasattr(draft, 'to_dict'):
            slides.append(draft.to_dict())
        elif isinstance(draft, dict):
            slides.append(draft)
        else:
            logger.warning(f"Unknown slide format: {type(draft)}")

    return create_presentation_v2(
        slides=slides,
        report_id=report_id,
        client_name=client_name,
        theme_id=theme_id
    )

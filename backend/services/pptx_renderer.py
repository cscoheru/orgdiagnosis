"""
PPTX Renderer Service

Renders structured SlideDraft objects to PowerPoint slides using python-pptx.
Supports multiple layouts and chart rendering.
"""

from typing import List, Dict, Any, Optional
from pathlib import Path
from loguru import logger
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class PPTXRenderer:
    """
    PPTX 渲染引擎

    将结构化的 SlideDraft 数据渲染为专业的 PowerPoint 演示文稿。
    支持多种布局:
    - title_slide: 封面
    - section_divider: 章节分隔页
    - bullet_points: 要点列表
    - two_columns: 双列对比
    - data_chart: 数据图表
    """

    # 布局映射
    LAYOUT_MAPPING = {
        "title_slide": 6,           # 空白布局
        "section_divider": 6,       # 空白布局
        "bullet_points": 6,         # 空白布局
        "two_columns": 6,           # 空白布局
        "swot_matrix": 6,
        "process_flow": 6,
        "five_dimensions_radar": 6,
        "gantt_chart": 6,
        "team_table": 6,
        "pricing_table": 6,
        "case_study": 6,
        "key_insight": 6,
    }

    # 颜色方案
    COLORS = {
        "primary": RGBColor(0, 82, 139),      # 蓝色
        "secondary": RGBColor(51, 102, 153),   # 深蓝
        "accent": RGBColor(231, 76, 60),       # 橙色
        "text_dark": RGBColor(51, 51, 51),     # 深灰
        "text_light": RGBColor(128, 128, 128), # 浅灰
        "background": RGBColor(245, 245, 245), # 浅灰背景
        "white": RGBColor(255, 255, 255),      # 白色
    }

    def __init__(
        self,
        template_path: Optional[str] = None,
        output_dir: str = "./output/pptx"
    ):
        """
        初始化 PPTX 渲染引擎

        Args:
            template_path: PPTX 模板文件路径
            output_dir: 输出目录
        """
        self.template_path = template_path
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        self.prs = None
        self._slide_width = Inches(13.333)
        self._slide_height = Inches(7.5)

        logger.info(f"PPTXRenderer initialized (output: {self.output_dir})")

    def render_report(
        self,
        slides: List[Dict[str, Any]],
        report_id: str,
        client_name: str
    ) -> str:
        """
        渲染完整报告

        Args:
            slides: SlideDraft 列表
            report_id: 报告ID
            client_name: 客户名称

        Returns:
            生成的 PPTX 文件路径
        """
        output_path = self.output_dir / f"{report_id}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"

        # 创建演示文稿
        if self.template_path and Path(self.template_path).exists():
            self.prs = Presentation(self.template_path)
        else:
            self.prs = Presentation()

        # 设置幻灯片尺寸 (16:9)
        self.prs.slide_width = self._slide_width
        self.prs.slide_height = self._slide_height

        logger.info(f"Creating presentation: {output_path}")
        logger.info(f"Total slides to render: {len(slides)}")

        try:
            for i, slide_data in enumerate(slides):
                self._render_slide(slide_data, i + 1)

            # 保存文件
            self.prs.save(output_path)
            logger.info(f"Report saved: {output_path}")

            return str(output_path)

        except Exception as e:
            logger.error(f"Error rendering report: {e}")
            raise

    def _render_slide(self, slide_data: Dict[str, Any], slide_index: int):
        """
        渲染单个幻灯片

        Args:
            slide_data: SlideDraft 数据
            slide_index: 幻灯片索引
        """
        layout = slide_data.get("layout", "bullet_points")
        layout_index = self.LAYOUT_MAPPING.get(layout, 6)  # 默认使用空白布局

        # 添加空白幻灯片
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)

        # 根据布局类型渲染内容
        if layout == "title_slide":
            self._render_title_slide(slide, slide_data)
        elif layout == "section_divider":
            self._render_section_divider(slide, slide_data)
        elif layout == "bullet_points":
            self._render_bullet_points(slide, slide_data)
        elif layout == "two_columns":
            self._render_two_columns(slide, slide_data)
        elif layout == "data_chart":
            self._render_data_chart(slide, slide_data)
        else:
            # 默认处理
            self._render_bullet_points(slide, slide_data)

    def _add_text_box(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_size: int = 18,
        bold: bool = False,
        color: RGBColor = None,
        alignment: PP_ALIGN = PP_ALIGN.LEFT
    ):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = alignment
        if color:
            p.font.color.rgb = color
        return txBox

    def _render_title_slide(self, slide, slide_data: Dict[str, Any]):
        """渲染封面页"""
        # 主标题
        title = slide_data.get("title", "咨询项目建议书")
        self._add_text_box(
            slide, 0.5, 2.5, 12.333, 1.5,
            title,
            font_size=44, bold=True,
            color=self.COLORS["primary"],
            alignment=PP_ALIGN.CENTER
        )

        # 客户名称
        client_name = slide_data.get("client_name", "")
        if client_name:
            self._add_text_box(
                slide, 0.5, 4.2, 12.333, 0.8,
                client_name,
                font_size=24, bold=False,
                color=self.COLORS["text_dark"],
                alignment=PP_ALIGN.CENTER
            )

        # 日期
        date_str = datetime.now().strftime("%Y年%m月")
        self._add_text_box(
            slide, 0.5, 5.5, 12.333, 0.5,
            date_str,
            font_size=16, bold=False,
            color=self.COLORS["text_light"],
            alignment=PP_ALIGN.CENTER
        )

    def _render_section_divider(self, slide, slide_data: Dict[str, Any]):
        """渲染章节分隔页"""
        # 章节标题
        subsection = slide_data.get("subsection", "")
        self._add_text_box(
            slide, 0.8, 2.8, 11.733, 1.2,
            subsection,
            font_size=36, bold=True,
            color=self.COLORS["primary"],
            alignment=PP_ALIGN.LEFT
        )

        # 章节描述
        title = slide_data.get("title", "")
        if title:
            self._add_text_box(
                slide, 0.8, 4.2, 11.733, 0.8,
                title,
                font_size=20, bold=False,
                color=self.COLORS["text_light"],
                alignment=PP_ALIGN.LEFT
            )

    def _render_bullet_points(self, slide, slide_data: Dict[str, Any]):
        """渲染要点列表页"""
        # 标题
        title = slide_data.get("title", "")
        self._add_text_box(
            slide, 0.5, 0.4, 12.333, 0.8,
            title,
            font_size=28, bold=True,
            color=self.COLORS["text_dark"],
            alignment=PP_ALIGN.LEFT
        )

        # 核心观点 (Action Title)
        key_message = slide_data.get("key_message", "")
        if key_message:
            self._add_text_box(
                slide, 0.5, 1.3, 12.333, 0.6,
                key_message,
                font_size=16, bold=True,
                color=self.COLORS["primary"],
                alignment=PP_ALIGN.LEFT
            )

        # 要点列表
        bullets = slide_data.get("bullets", [])
        if bullets:
            bullet_text = "\n\n".join([f"• {bullet}" for bullet in bullets])
            self._add_text_box(
                slide, 0.5, 2.2, 12.333, 4.5,
                bullet_text,
                font_size=18, bold=False,
                color=self.COLORS["text_dark"],
                alignment=PP_ALIGN.LEFT
            )

        # 证据来源
        source_ref = slide_data.get("source_ref", "")
        if source_ref:
            self._add_text_box(
                slide, 0.5, 6.8, 12.333, 0.4,
                f"来源: {source_ref}",
                font_size=10, bold=False,
                color=self.COLORS["text_light"],
                alignment=PP_ALIGN.LEFT
            )

    def _render_two_columns(self, slide, slide_data: Dict[str, Any]):
        """渲染双列对比页"""
        # 标题
        title = slide_data.get("title", "")
        self._add_text_box(
            slide, 0.5, 0.4, 12.333, 0.8,
            title,
            font_size=28, bold=True,
            color=self.COLORS["text_dark"],
            alignment=PP_ALIGN.CENTER
        )

        # 左右两列内容
        bullets = slide_data.get("bullets", [])
        if bullets:
            mid = len(bullets) // 2
            left_bullets = bullets[:mid]
            right_bullets = bullets[mid:]

            # 左列
            if left_bullets:
                left_text = "\n\n".join([f"• {b}" for b in left_bullets])
                self._add_text_box(
                    slide, 0.5, 1.5, 5.9, 5.5,
                    left_text,
                    font_size=16, bold=False,
                    color=self.COLORS["text_dark"],
                    alignment=PP_ALIGN.LEFT
                )

            # 右列
            if right_bullets:
                right_text = "\n\n".join([f"• {b}" for b in right_bullets])
                self._add_text_box(
                    slide, 6.9, 1.5, 5.9, 5.5,
                    right_text,
                    font_size=16, bold=False,
                    color=self.COLORS["text_dark"],
                    alignment=PP_ALIGN.LEFT
                )

    def _render_data_chart(self, slide, slide_data: Dict[str, Any]):
        """渲染数据图表页"""
        # 标题
        title = slide_data.get("title", "")
        self._add_text_box(
            slide, 0.5, 0.4, 12.333, 0.8,
            title,
            font_size=28, bold=True,
            color=self.COLORS["text_dark"],
            alignment=PP_ALIGN.CENTER
        )

        # 图表数据 - 渲染为 PPTX 表格
        chart_data = slide_data.get("chart_data", {})
        rows = slide_data.get("chart_rows", [])

        if rows and isinstance(rows[0], dict):
            # 结构化表格数据: [{"label": "...", "values": [...]}]
            self._render_table_from_rows(slide, rows, 1.5)
        elif chart_data and isinstance(chart_data, dict):
            # 键值对数据: {"指标1": 85, "指标2": 72}
            self._render_table_from_kv(slide, chart_data, 1.5)
        else:
            # Fallback: 渲染 bullets
            bullets = slide_data.get("bullets", [])
            if bullets:
                bullet_text = "\n\n".join([f"• {b}" for b in bullets])
                self._add_text_box(
                    slide, 0.5, 1.5, 12.333, 5.0,
                    bullet_text,
                    font_size=16, bold=False,
                    color=self.COLORS["text_dark"],
                    alignment=PP_ALIGN.LEFT
                )

    def _render_table_from_rows(self, slide, rows: list, top: float):
        """从结构化行数据渲染表格"""
        if not rows:
            return

        # 确定列名
        first_row = rows[0]
        if isinstance(first_row, dict) and "values" in first_row:
            # rows = [{"label": "X", "values": [1, 2, 3]}, ...]
            headers = first_row.get("columns", [f"指标{i+1}" for i in range(len(first_row.get("values", [])))])
            num_cols = len(headers)
            table_data = [[first_row.get("label", "")] + [str(v) for v in first_row.get("values", [])]]
            for row in rows[1:]:
                table_data.append([row.get("label", "")] + [str(v) for v in row.get("values", [])])
        else:
            headers = list(first_row.keys())
            num_cols = len(headers)
            table_data = [[str(first_row.get(h, "")) for h in headers]]
            for row in rows[1:]:
                table_data.append([str(row.get(h, "")) for h in headers])

        if not table_data:
            return

        num_rows = len(table_data)
        col_width = min(2.5, 12.0 / max(num_cols, 1))

        try:
            table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(top), Inches(col_width * num_cols), Inches(0.4 * num_rows))
            table = table_shape.table

            # Header row styling
            for j, header in enumerate(headers):
                cell = table.cell(0, j)
                cell.text = header
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = self.COLORS["white"]
                # Blue background
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = self.COLORS["primary"]

            # Data rows
            for i in range(1, num_rows):
                for j in range(num_cols):
                    cell = table.cell(i, j)
                    cell.text = table_data[i][j] if j < len(table_data[i]) else ""
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(11)
                    # Alternating row colors
                    if i % 2 == 0:
                        cell_fill = cell.fill
                        cell_fill.solid()
                        cell_fill.fore_color.rgb = self.COLORS["background"]
        except Exception as e:
            logger.warning(f"Failed to render table: {e}")

    def _render_table_from_kv(self, slide, chart_data: dict, top: float):
        """从键值对渲染表格"""
        if not chart_data:
            return

        items = list(chart_data.items())
        num_rows = len(items) + 1  # +1 for header
        num_cols = 2

        try:
            table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(top), Inches(12.0), Inches(0.4 * num_rows))
            table = table_shape.table

            # Header
            for j, header in enumerate(["指标", "数值"]):
                cell = table.cell(0, j)
                cell.text = header
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = self.COLORS["white"]
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = self.COLORS["primary"]

            # Data rows
            for i, (key, value) in enumerate(items):
                cell_key = table.cell(i + 1, 0)
                cell_key.text = str(key)
                for p in cell_key.text_frame.paragraphs:
                    p.font.size = Pt(11)

                cell_val = table.cell(i + 1, 1)
                cell_val.text = str(value)
                for p in cell_val.text_frame.paragraphs:
                    p.font.size = Pt(11)

                if i % 2 == 0:
                    for cell in [cell_key, cell_val]:
                        fill = cell.fill
                        fill.solid()
                        fill.fore_color.rgb = self.COLORS["background"]
        except Exception as e:
            logger.warning(f"Failed to render KV table: {e}")


# 便捷函数
def create_presentation(
    slides: List[Dict[str, Any]],
    report_id: str,
    client_name: str,
    output_dir: str = "./output/pptx",
    template_id: Optional[str] = None,
) -> str:
    """
    创建演示文稿的便捷函数

    Args:
        slides: SlideDraft 列表
        report_id: 报告ID
        client_name: 客户名称
        output_dir: 输出目录
        template_id: 模板ID (用于未来模板支持)

    Returns:
        生成的 PPTX 文件路径
    """
    logger.info(f"Creating presentation with template_id: {template_id}")
    renderer = PPTXRenderer(output_dir=output_dir)
    return renderer.render_report(slides, report_id, client_name)

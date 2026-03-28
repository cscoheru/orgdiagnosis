"""
PPT 报告生成器 — 模板占位符替换

刚性约束:
1. 打开现有模板文件
2. 查找 {{ PLACEHOLDER }} 占位符
3. 替换时保留原有字体颜色、大小和段落样式
"""
import os
import re
from typing import Any
from pptx import Presentation


def generate_org_ppt(
    template_path: str,
    data: dict[str, Any],
    output_path: str,
) -> dict[str, Any]:
    """生成组织架构 PPT 报告"""
    try:
        if not os.path.exists(template_path):
            return {"success": False, "error": f"模板文件不存在: {template_path}"}

        prs = Presentation(template_path)
        placeholders_replaced = 0
        placeholder_pattern = re.compile(r"\{\{\s*(\w+)\s*\}\}")

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    replaced = _process_text_frame(shape.text_frame, data, placeholder_pattern)
                    placeholders_replaced += replaced

                if shape.has_table:
                    replaced = _process_table(shape.table, data, placeholder_pattern)
                    placeholders_replaced += replaced

        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)

        prs.save(output_path)

        return {
            "success": True,
            "placeholders_replaced": placeholders_replaced,
            "output_path": output_path,
            "file_size": os.path.getsize(output_path),
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


def _process_text_frame(
    text_frame: Any,
    data: dict[str, Any],
    placeholder_pattern: re.Pattern,
) -> int:
    """处理文本框中的占位符"""
    replaced_count = 0

    for paragraph in text_frame.paragraphs:
        full_text = paragraph.text
        if not full_text:
            continue

        matches = list(placeholder_pattern.finditer(full_text))
        if not matches:
            continue

        first_run_style = {}
        if paragraph.runs:
            first_run_style = _extract_run_style(paragraph.runs[0])

        new_text = full_text
        for match in matches:
            placeholder_key = match.group(1)
            if placeholder_key in data:
                replacement_value = str(data[placeholder_key])
                new_text = new_text.replace(match.group(0), replacement_value)
                replaced_count += 1

        for run in paragraph.runs:
            run.text = ""

        if paragraph.runs:
            paragraph.runs[0].text = new_text
            _apply_run_style(paragraph.runs[0], first_run_style)
        else:
            run = paragraph.add_run()
            run.text = new_text
            _apply_run_style(run, first_run_style)

    return replaced_count


def _extract_run_style(run: Any) -> dict[str, Any]:
    """提取 run 的样式属性"""
    style = {}
    if run.font:
        font = run.font
        style["font_name"] = font.name
        style["font_size"] = font.size
        style["font_bold"] = font.bold
        style["font_italic"] = font.italic
        style["font_underline"] = font.underline
        try:
            if font.color and font.color.rgb:
                style["font_color"] = font.color.rgb
        except (AttributeError, TypeError):
            pass
        if hasattr(font, "strike"):
            style["font_strike"] = font.strike
    return style


def _apply_run_style(run: Any, style: dict[str, Any]) -> None:
    """将样式应用到 run"""
    if not style:
        return

    if run.font:
        font = run.font
        if style.get("font_name"):
            font.name = style["font_name"]
        if style.get("font_size"):
            font.size = style["font_size"]
        if style.get("font_bold") is not None:
            font.bold = style["font_bold"]
        if style.get("font_italic") is not None:
            font.italic = style["font_italic"]
        if style.get("font_underline") is not None:
            font.underline = style["font_underline"]
        try:
            if style.get("font_color"):
                font.color.rgb = style["font_color"]
        except (AttributeError, TypeError):
            pass
        if style.get("font_strike") is not None:
            font.strike = style["font_strike"]


def _process_table(
    table,
    data: dict[str, Any],
    placeholder_pattern: re.Pattern,
) -> int:
    """处理表格中的占位符"""
    replaced_count = 0
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                replaced_count += _process_text_frame(
                    cell.text_frame, data, placeholder_pattern
                )
    return replaced_count

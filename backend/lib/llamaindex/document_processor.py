"""
Consulting Document Processor

Handles loading and parsing of consulting documents (PDF, DOCX, PPTX).
Extracts metadata and prepares documents for indexing.
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime
from loguru import logger

from llama_index.core import Document
from llama_index.core.schema import MetadataMode


class ConsultingDocumentProcessor:
    """
    咨询文档处理器

    支持 PDF/DOCX/PPTX/MD 格式的咨询报告解析
    """

    # 支持的文件扩展名
    SUPPORTED_EXTENSIONS = [".pdf", ".docx", ".pptx", ".md", ".txt"]

    # 咨询领域关键词映射
    CATEGORY_KEYWORDS = {
        "strategy": ["战略", "战略规划", "市场洞察", "业务设计", "竞争分析", "愿景", "使命", "商业模式"],
        "hr": ["人力资源", "组织架构", "人才盘点", "招聘", "培养", "离职", "敬业度"],
        "performance": ["绩效", "考核", "KPI", "OKR", "目标管理", "指标分解", "绩效辅导"],
        "compensation": ["薪酬", "激励", "工资", "奖金", "股权", "福利", "固浮比", "薪酬体系"],
        "finance": ["财务", "预算", "成本", "资金", "投资", "现金流"],
        "operations": ["运营", "流程", "供应链", "生产", "质量", "数字化"],
    }

    def __init__(self, data_dir: str = "./data/historical_reports"):
        """
        初始化文档处理器

        Args:
            data_dir: 历史报告存放目录
        """
        self.data_dir = Path(data_dir)
        self.data_dir.mkdir(parents=True, exist_ok=True)

        # 确保子目录存在
        for category in ["strategy", "hr", "performance", "compensation", "finance", "operations"]:
            (self.data_dir / category).mkdir(exist_ok=True)

    def load_documents(
        self,
        category: Optional[str] = None,
        file_paths: Optional[List[str]] = None
    ) -> List[Document]:
        """
        加载文档

        Args:
            category: 指定加载的类别 (strategy/hr/finance/operations)
            file_paths: 指定加载的文件路径列表

        Returns:
            Document 列表
        """
        documents = []

        if file_paths:
            # 加载指定文件
            for file_path in file_paths:
                doc = self._load_single_file(file_path)
                if doc:
                    documents.extend(doc)
        elif category:
            # 加载指定类别的所有文件
            category_dir = self.data_dir / category
            if category_dir.exists():
                documents = self._load_directory(category_dir, category)
        else:
            # 加载所有类别的文件
            for cat in ["strategy", "hr", "performance", "compensation", "finance", "operations"]:
                cat_dir = self.data_dir / cat
                if cat_dir.exists():
                    documents.extend(self._load_directory(cat_dir, cat))

        # 清理元数据，确保 ChromaDB 兼容
        documents = [self._sanitize_metadata(doc) for doc in documents]

        logger.info(f"Loaded {len(documents)} document chunks")
        return documents

    def _load_directory(self, directory: Path, category: str) -> List[Document]:
        """加载目录下的所有文档"""
        documents = []

        for file_path in directory.rglob("*"):
            if file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                docs = self._load_single_file(str(file_path), category)
                if docs:
                    documents.extend(docs)

        return documents

    def _load_single_file(
        self,
        file_path: str,
        category: Optional[str] = None
    ) -> List[Document]:
        """
        加载单个文件

        Args:
            file_path: 文件路径
            category: 文档类别

        Returns:
            Document 列表 (一个文件可能被分成多个 chunk)
        """
        path = Path(file_path)

        if not path.exists():
            logger.warning(f"File not found: {file_path}")
            return []

        if path.suffix.lower() not in self.SUPPORTED_EXTENSIONS:
            logger.warning(f"Unsupported file type: {path.suffix}")
            return []

        try:
            if path.suffix.lower() == ".pdf":
                return self._load_pdf(file_path, category)
            elif path.suffix.lower() == ".docx":
                return self._load_docx(file_path, category)
            elif path.suffix.lower() == ".pptx":
                return self._load_pptx(file_path, category)
            elif path.suffix.lower() in [".md", ".txt"]:
                return self._load_text(file_path, category)
        except Exception as e:
            logger.error(f"Error loading {file_path}: {e}")
            return []

        return []

    def _load_pdf(self, file_path: str, category: Optional[str]) -> List[Document]:
        """加载 PDF 文件"""
        try:
            from llama_index.readers.file import PyMuPDFReader

            reader = PyMuPDFReader()
            docs = reader.load_data(Path(file_path))

            # 添加元数据
            for doc in docs:
                doc.metadata.update(self._extract_metadata(file_path, category, "pdf"))

            return docs
        except Exception as e:
            logger.warning(f"PyMuPDFReader failed: {e}, using fallback")
            return self._load_pdf_fallback(file_path, category)

    def _load_pdf_fallback(self, file_path: str, category: Optional[str]) -> List[Document]:
        """PDF 加载备用方案"""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()

            metadata = self._extract_metadata(file_path, category, "pdf")
            return [Document(text=text, metadata=metadata)]
        except ImportError:
            logger.error("PyMuPDF not installed")
            return []

    def _load_docx(self, file_path: str, category: Optional[str]) -> List[Document]:
        """加载 DOCX 文件"""
        try:
            from llama_index.readers.file import DocxReader

            reader = DocxReader()
            docs = reader.load_data(file=Path(file_path))

            for doc in docs:
                doc.metadata.update(self._extract_metadata(file_path, category, "docx"))

            return docs
        except ImportError:
            logger.warning("DocxReader not available, using fallback")
            return self._load_docx_fallback(file_path, category)

    def _load_docx_fallback(self, file_path: str, category: Optional[str]) -> List[Document]:
        """DOCX 加载备用方案"""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])

            metadata = self._extract_metadata(file_path, category, "docx")
            return [Document(text=text, metadata=metadata)]
        except ImportError:
            logger.error("python-docx not installed")
            return []

    def _load_pptx(self, file_path: str, category: Optional[str]) -> List[Document]:
        """加载 PPTX 文件"""
        try:
            from llama_index.readers.file import PptxReader

            reader = PptxReader()
            docs = reader.load_data(file=Path(file_path))

            for doc in docs:
                doc.metadata.update(self._extract_metadata(file_path, category, "pptx"))

            return docs
        except ImportError:
            logger.warning("PptxReader not available, using fallback")
            return self._load_pptx_fallback(file_path, category)

    def _load_pptx_fallback(self, file_path: str, category: Optional[str]) -> List[Document]:
        """PPTX 加载备用方案"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text = ""

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

            metadata = self._extract_metadata(file_path, category, "pptx")
            return [Document(text=text, metadata=metadata)]
        except ImportError:
            logger.error("python-pptx not installed")
            return []

    def _load_text(self, file_path: str, category: Optional[str]) -> List[Document]:
        """加载文本文件 (MD/TXT)"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()

            metadata = self._extract_metadata(file_path, category, "text")
            return [Document(text=text, metadata=metadata)]
        except Exception as e:
            logger.error(f"Error reading text file: {e}")
            return []

    def _extract_metadata(
        self,
        file_path: str,
        category: Optional[str],
        file_type: str
    ) -> Dict[str, Any]:
        """
        提取文档元数据

        Args:
            file_path: 文件路径
            category: 文档类别
            file_type: 文件类型

        Returns:
            元数据字典
        """
        path = Path(file_path)

        # 从路径推断类别
        inferred_category = category
        if not inferred_category:
            for cat in ["strategy", "hr", "finance", "operations"]:
                if cat in str(path).lower():
                    inferred_category = cat
                    break

        metadata = {
            "file_name": path.name,
            "file_path": str(path),
            "file_type": file_type,
            "category": inferred_category or "general",
            "created_at": datetime.now().isoformat(),
        }

        return metadata

    def extract_keywords(self, text: str) -> List[str]:
        """
        从文本中提取关键词

        用于自动分类和增强检索
        """
        keywords = []

        for category, kw_list in self.CATEGORY_KEYWORDS.items():
            for kw in kw_list:
                if kw in text:
                    keywords.append(kw)

        return list(set(keywords))

    def _sanitize_metadata(self, doc: Document) -> Document:
        """
        清理文档元数据，确保所有值都是 ChromaDB 支持的类型

        ChromaDB 只支持 str, int, float, None
        """
        sanitized = {}
        for key, value in doc.metadata.items():
            if value is None:
                sanitized[key] = None
            elif isinstance(value, (str, int, float)):
                sanitized[key] = value
            elif isinstance(value, (list, dict)):
                # 将列表和字典转换为字符串
                sanitized[key] = str(value)
            else:
                sanitized[key] = str(value)

        doc.metadata = sanitized
        return doc

    def infer_category_from_content(self, text: str) -> str:
        """
        根据内容推断文档类别

        Returns:
            推断的类别: strategy/hr/finance/operations/general
        """
        scores = {cat: 0 for cat in self.CATEGORY_KEYWORDS.keys()}

        for category, keywords in self.CATEGORY_KEYWORDS.items():
            for kw in keywords:
                if kw in text:
                    scores[category] += 1

        if max(scores.values()) == 0:
            return "general"

        return max(scores, key=scores.get)


# 便捷函数
def load_consulting_documents(
    data_dir: str = "./data/historical_reports",
    category: Optional[str] = None
) -> List[Document]:
    """
    加载咨询文档的便捷函数

    Args:
        data_dir: 历史报告目录
        category: 指定类别

    Returns:
        Document 列表
    """
    processor = ConsultingDocumentProcessor(data_dir)
    return processor.load_documents(category=category)

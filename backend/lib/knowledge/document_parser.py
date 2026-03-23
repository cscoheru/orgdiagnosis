"""
文档解析器

支持解析:
- PPTX: python-pptx
- PDF: pdfplumber

按页拆分，提取文字和表格内容

Created: 2026-03-22
"""

import os
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from loguru import logger


@dataclass
class ParsedPage:
    """解析后的页面"""
    page_number: int
    content: str  # 纯文本内容
    sections: List[Dict[str, Any]] = field(default_factory=list)  # 结构化内容块
    has_table: bool = False
    has_image: bool = False
    word_count: int = 0


@dataclass
class ParsedDocument:
    """解析后的文档"""
    filename: str
    file_type: str
    page_count: int
    pages: List[ParsedPage]
    title: Optional[str] = None
    author: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "filename": self.filename,
            "file_type": self.file_type,
            "page_count": self.page_count,
            "title": self.title,
            "author": self.author,
            "metadata": self.metadata,
            "pages": [
                {
                    "page_number": p.page_number,
                    "content": p.content,
                    "sections": p.sections,
                    "has_table": p.has_table,
                    "has_image": p.has_image,
                    "word_count": p.word_count
                }
                for p in self.pages
            ]
        }


class DocumentParser:
    """
    文档解析器

    支持:
    - PPTX (python-pptx)
    - PDF (pdfplumber)
    - DOCX (python-docx)
    - XLSX/XLS (openpyxl)
    - MD (markdown, plain text)
    - JSON (text extraction)
    - Images (PNG, JPG, JPEG - metadata only for now)
    """

    SUPPORTED_EXTENSIONS = [".pptx", ".pdf", ".docx", ".xlsx", ".xls", ".md", ".json", ".png", ".jpg", ".jpeg"]

    def __init__(self):
        """初始化解析器"""
        self._check_dependencies()

    def _check_dependencies(self):
        """检查依赖库"""
        try:
            from pptx import Presentation
            self._pptx_available = True
        except ImportError:
            self._pptx_available = False
            logger.warning("python-pptx not installed, PPTX parsing disabled")

        try:
            import pdfplumber
            self._pdfplumber_available = True
        except ImportError:
            self._pdfplumber_available = False
            logger.warning("pdfplumber not installed, PDF parsing disabled")

        try:
            from docx import Document
            self._docx_available = True
        except ImportError:
            self._docx_available = False
            logger.warning("python-docx not installed, DOCX parsing disabled")

        try:
            from openpyxl import load_workbook
            self._xlsx_available = True
        except ImportError:
            self._xlsx_available = False
            logger.warning("openpyxl not installed, XLSX parsing disabled")

        try:
            from PIL import Image
            self._pil_available = True
        except ImportError:
            self._pil_available = False
            logger.warning("Pillow not installed, image parsing limited")

    def parse(self, file_path: str) -> ParsedDocument:
        """
        解析文档

        Args:
            file_path: 文件路径

        Returns:
            ParsedDocument 对象
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        if ext == ".pptx":
            return self._parse_pptx(file_path)
        elif ext == ".pdf":
            return self._parse_pdf(file_path)
        elif ext == ".docx":
            return self._parse_docx(file_path)
        elif ext in [".xlsx", ".xls"]:
            return self._parse_xlsx(file_path)
        elif ext == ".md":
            return self._parse_markdown(file_path)
        elif ext == ".json":
            return self._parse_json(file_path)
        elif ext in [".png", ".jpg", ".jpeg"]:
            return self._parse_image(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _parse_pptx(self, file_path: str) -> ParsedDocument:
        """解析 PPTX 文件"""
        if not self._pptx_available:
            raise ImportError("python-pptx not installed")

        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        path = Path(file_path)
        prs = Presentation(file_path)
        pages = []

        # 提取元数据
        doc_title = None
        doc_author = None
        core_props = prs.core_properties
        if core_props:
            doc_title = core_props.title or path.stem
            doc_author = core_props.author

        for slide_idx, slide in enumerate(prs.slides, start=1):
            sections = []
            all_text = []
            has_table = False
            has_image = False

            # 遍历所有形状
            for shape in slide.shapes:
                # 文本框
                if shape.has_text_frame:
                    text = self._extract_text_frame(shape.text_frame)
                    if text.strip():
                        sections.append({
                            "type": "text",
                            "content": text,
                            "shape_name": getattr(shape, 'name', None)
                        })
                        all_text.append(text)

                # 表格
                elif shape.has_table:
                    has_table = True
                    table_content = self._extract_pptx_table(shape.table)
                    sections.append({
                        "type": "table",
                        "content": table_content,
                        "rows": len(shape.table.rows),
                        "cols": len(shape.table.columns)
                    })
                    all_text.append(table_content)

                # 图片
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True
                    sections.append({
                        "type": "image",
                        "shape_name": getattr(shape, 'name', None)
                    })

            # 提取备注
            if slide.has_notes_slide:
                notes_text = self._extract_notes(slide.notes_slide)
                if notes_text.strip():
                    sections.append({
                        "type": "notes",
                        "content": notes_text
                    })
                    all_text.append(notes_text)

            content = "\n\n".join(all_text)
            pages.append(ParsedPage(
                page_number=slide_idx,
                content=content,
                sections=sections,
                has_table=has_table,
                has_image=has_image,
                word_count=len(content)
            ))

        logger.info(f"PPTX parsed: {file_path}, {len(pages)} slides")

        return ParsedDocument(
            filename=path.name,
            file_type="pptx",
            page_count=len(pages),
            pages=pages,
            title=doc_title,
            author=doc_author,
            metadata={"source": "python-pptx"}
        )

    def _parse_pdf(self, file_path: str) -> ParsedDocument:
        """解析 PDF 文件"""
        if not self._pdfplumber_available:
            raise ImportError("pdfplumber not installed")

        import pdfplumber

        path = Path(file_path)
        pages = []

        with pdfplumber.open(file_path) as pdf:
            # 提取元数据
            doc_title = path.stem
            doc_author = None
            if pdf.metadata:
                doc_title = pdf.metadata.get('Title') or path.stem
                doc_author = pdf.metadata.get('Author')

            for page_idx, page in enumerate(pdf.pages, start=1):
                sections = []
                all_text = []
                has_table = False

                # 提取文本
                text = page.extract_text() or ""
                if text.strip():
                    sections.append({
                        "type": "text",
                        "content": text.strip()
                    })
                    all_text.append(text.strip())

                # 提取表格
                tables = page.extract_tables()
                if tables:
                    has_table = True
                    for table_idx, table in enumerate(tables):
                        table_text = self._format_pdf_table(table)
                        sections.append({
                            "type": "table",
                            "content": table_text,
                            "table_index": table_idx
                        })
                        all_text.append(table_text)

                content = "\n\n".join(all_text)
                pages.append(ParsedPage(
                    page_number=page_idx,
                    content=content,
                    sections=sections,
                    has_table=has_table,
                    has_image=False,  # pdfplumber 不处理图片
                    word_count=len(content)
                ))

        logger.info(f"PDF parsed: {file_path}, {len(pages)} pages")

        return ParsedDocument(
            filename=path.name,
            file_type="pdf",
            page_count=len(pages),
            pages=pages,
            title=doc_title,
            author=doc_author,
            metadata={"source": "pdfplumber"}
        )

    def _extract_text_frame(self, text_frame) -> str:
        """提取文本框内容"""
        paragraphs = []
        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                paragraphs.append(text)
        return "\n".join(paragraphs)

    def _extract_pptx_table(self, table) -> str:
        """提取 PPTX 表格内容"""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                cells.append(cell_text)
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _extract_notes(self, notes_slide) -> str:
        """提取备注内容"""
        texts = []
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text.strip())
        return "\n".join(texts)

    def _format_pdf_table(self, table: List[List[str]]) -> str:
        """格式化 PDF 表格"""
        rows = []
        for row in table:
            if row:  # 跳过空行
                cells = [str(cell) if cell else "" for cell in row]
                rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _parse_docx(self, file_path: str) -> ParsedDocument:
        """解析 DOCX 文件"""
        if not self._docx_available:
            raise ImportError("python-docx not installed")

        from docx import Document

        path = Path(file_path)
        doc = Document(file_path)
        pages = []

        # 提取元数据
        doc_title = None
        doc_author = None
        core_props = doc.core_properties
        if core_props:
            doc_title = core_props.title or path.stem
            doc_author = core_props.author

        # 提取所有段落
        all_text = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                all_text.append(text)

        # 提取表格
        for table in doc.tables:
            table_text = self._extract_docx_table(table)
            all_text.append(table_text)

        content = "\n\n".join(all_text)
        pages.append(ParsedPage(
            page_number=1,
            content=content,
            sections=[{"type": "text", "content": content}],
            has_table=len(doc.tables) > 0,
            has_image=False,
            word_count=len(content)
        ))

        logger.info(f"DOCX parsed: {file_path}")

        return ParsedDocument(
            filename=path.name,
            file_type="docx",
            page_count=1,
            pages=pages,
            title=doc_title,
            author=doc_author,
            metadata={"source": "python-docx"}
        )

    def _extract_docx_table(self, table) -> str:
        """提取 DOCX 表格内容"""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                cells.append(cell_text)
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _parse_xlsx(self, file_path: str) -> ParsedDocument:
        """解析 XLSX/XLS 文件"""
        if not self._xlsx_available:
            raise ImportError("openpyxl not installed")

        from openpyxl import load_workbook

        path = Path(file_path)
        wb = load_workbook(file_path, data_only=True)
        pages = []

        # 每个工作表作为一个页面
        for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
            sheet = wb[sheet_name]
            sections = []
            all_text = []

            # 提取所有单元格
            rows_data = []
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_text):  # 跳过空行
                    rows_data.append(" | ".join(row_text))

            if rows_data:
                content = "\n".join(rows_data)
                sections.append({
                    "type": "spreadsheet",
                    "sheet_name": sheet_name,
                    "rows": len(rows_data)
                })
                all_text.append(f"## {sheet_name}\n{content}")

                pages.append(ParsedPage(
                    page_number=sheet_idx,
                    content="\n".join(all_text),
                    sections=sections,
                    has_table=True,
                    has_image=False,
                    word_count=len("\n".join(all_text))
                ))

        logger.info(f"XLSX parsed: {file_path}, {len(pages)} sheets")

        return ParsedDocument(
            filename=path.name,
            file_type="xlsx",
            page_count=len(pages),
            pages=pages,
            title=path.stem,
            author=None,
            metadata={"source": "openpyxl", "sheets": wb.sheetnames}
        )

    def _parse_markdown(self, file_path: str) -> ParsedDocument:
        """解析 Markdown 文件"""
        path = Path(file_path)

        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # 提取标题作为文档标题
        lines = content.split('\n')
        doc_title = path.stem
        for line in lines:
            if line.startswith('# '):
                doc_title = line[2:].strip()
                break

        pages = [ParsedPage(
            page_number=1,
            content=content,
            sections=[{"type": "markdown", "content": content}],
            has_table=False,
            has_image=False,
            word_count=len(content)
        )]

        logger.info(f"Markdown parsed: {file_path}")

        return ParsedDocument(
            filename=path.name,
            file_type="md",
            page_count=1,
            pages=pages,
            title=doc_title,
            author=None,
            metadata={"source": "plain_text"}
        )

    def _parse_json(self, file_path: str) -> ParsedDocument:
        """解析 JSON 文件"""
        import json

        path = Path(file_path)

        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        # 格式化 JSON 为可读文本
        content = json.dumps(data, ensure_ascii=False, indent=2)

        pages = [ParsedPage(
            page_number=1,
            content=content,
            sections=[{"type": "json", "content": content}],
            has_table=False,
            has_image=False,
            word_count=len(content)
        )]

        logger.info(f"JSON parsed: {file_path}")

        return ParsedDocument(
            filename=path.name,
            file_type="json",
            page_count=1,
            pages=pages,
            title=path.stem,
            author=None,
            metadata={"source": "json_parser"}
        )

    def _parse_image(self, file_path: str) -> ParsedDocument:
        """解析图片文件"""
        path = Path(file_path)
        ext = path.suffix.lower()

        content = f"[图片文件: {path.name}]"

        # 如果有 PIL，提取图片信息
        if self._pil_available:
            try:
                from PIL import Image
                with Image.open(file_path) as img:
                    width, height = img.size
                    format_name = img.format or "Unknown"
                    mode = img.mode
                    content = f"[图片文件: {path.name}]\n格式: {format_name}\n尺寸: {width}x{height}\n模式: {mode}"
            except Exception as e:
                logger.warning(f"Failed to read image info: {e}")

        pages = [ParsedPage(
            page_number=1,
            content=content,
            sections=[{"type": "image", "filename": path.name}],
            has_table=False,
            has_image=True,
            word_count=0
        )]

        logger.info(f"Image parsed: {file_path}")

        return ParsedDocument(
            filename=path.name,
            file_type=ext[1:],  # 去掉点
            page_count=1,
            pages=pages,
            title=path.stem,
            author=None,
            metadata={"source": "image_parser"}
        )

    def get_text_preview(self, parsed: ParsedDocument, max_chars: int = 500) -> str:
        """
        获取文档文本预览

        合并前几页的内容，用于分类
        """
        texts = []
        total_chars = 0

        for page in parsed.pages[:5]:  # 只取前5页
            if total_chars + len(page.content) > max_chars:
                texts.append(page.content[:max_chars - total_chars])
                break
            texts.append(page.content)
            total_chars += len(page.content)

        return "\n\n".join(texts)

    def get_full_text(self, parsed: ParsedDocument) -> str:
        """获取文档完整文本"""
        return "\n\n".join([p.content for p in parsed.pages])


# 便捷函数
def parse_document(file_path: str) -> ParsedDocument:
    """解析文档"""
    parser = DocumentParser()
    return parser.parse(file_path)


def parse_and_store(file_path: str, store, project_id: str = None) -> str:
    """
    解析文档并存储到知识库

    Args:
        file_path: 文件路径
        store: KnowledgeBaseStore 实例
        project_id: 项目ID

    Returns:
        文档ID
    """
    import uuid
    from pathlib import Path

    parser = DocumentParser()
    parsed = parser.parse(file_path)
    path = Path(file_path)

    # 创建文档记录
    doc_id = str(uuid.uuid4())
    store.create_document({
        "id": doc_id,
        "filename": parsed.filename,
        "file_type": parsed.file_type,
        "file_path": file_path,
        "file_size": path.stat().st_size if path.exists() else None,
        "project_id": project_id,
        "title": parsed.title,
        "author": parsed.author,
        "page_count": parsed.page_count,
        "metadata": parsed.metadata
    })

    # 创建页面记录
    for page in parsed.pages:
        store.create_page({
            "id": str(uuid.uuid4()),
            "document_id": doc_id,
            "page_number": page.page_number,
            "content": page.content,
            "sections": page.sections
        })

    logger.info(f"Document stored: {doc_id}, {parsed.page_count} pages")

    return doc_id

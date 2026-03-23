"""
PPTX Template Analyzer

Parses uploaded PowerPoint templates and extracts layout information including:
- Slide layout names and IDs
- Placeholder positions and types
- Suggested element counts
- Semantic descriptions (via AI)

Created: 2026-03-21
"""

import os
import json
import logging
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, asdict
from pathlib import Path

logger = logging.getLogger(__name__)

# Try to import python-pptx
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. Template analysis will be limited.")


@dataclass
class PlaceholderInfo:
    """Information about a placeholder in a layout"""
    placeholder_id: int
    placeholder_type: str  # title, body, picture, chart, table
    position: Tuple[float, float, float, float]  # left, top, width, height
    name: str = ""


@dataclass
class LayoutManifest:
    """Manifest for a single layout extracted from a template"""
    layout_id: str
    layout_name: str
    source_template: str
    description: str = ""
    category: str = "PARALLEL"  # MATRIX, PROCESS, PARALLEL, TABLE, TIMELINE, DATA_VIZ, KEY_INSIGHT
    element_count_range: Tuple[int, int] = (1, 6)
    placeholders: List[Dict[str, Any]] = None
    keywords: List[str] = None
    thumbnail_path: str = ""

    def __post_init__(self):
        if self.placeholders is None:
            self.placeholders = []
        if self.keywords is None:
            self.keywords = []

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary"""
        return asdict(self)


class TemplateAnalyzer:
    """
    Analyzes PPTX templates to extract layout information.

    Usage:
        analyzer = TemplateAnalyzer()
        manifests = analyzer.analyze_template("template.pptx")
    """

    def __init__(self, use_ai_description: bool = True):
        """
        Initialize the analyzer.

        Args:
            use_ai_description: Whether to use AI for generating semantic descriptions
        """
        self.use_ai_description = use_ai_description
        self._ai_service = None

    def analyze_template(self, pptx_path: str) -> List[LayoutManifest]:
        """
        Analyze a PPTX template and extract layout manifests.

        Args:
            pptx_path: Path to the PPTX template file

        Returns:
            List of LayoutManifest for each slide layout in the template
        """
        if not PPTX_AVAILABLE:
            logger.error("python-pptx not available. Cannot analyze template.")
            return []

        if not os.path.exists(pptx_path):
            logger.error(f"Template file not found: {pptx_path}")
            return []

        try:
            prs = Presentation(pptx_path)
            manifests = []
            template_name = Path(pptx_path).stem

            logger.info(f"Analyzing template: {template_name} with {len(prs.slide_layouts)} layouts")

            for idx, layout in enumerate(prs.slide_layouts):
                manifest = self._analyze_layout(layout, idx, template_name)
                if manifest:
                    manifests.append(manifest)

            # Generate AI descriptions if enabled
            if self.use_ai_description and manifests:
                manifests = self._generate_ai_descriptions(manifests)

            logger.info(f"Extracted {len(manifests)} layouts from template")
            return manifests

        except Exception as e:
            logger.error(f"Failed to analyze template: {e}")
            return []

    def _analyze_layout(
        self,
        layout,
        index: int,
        template_name: str
    ) -> Optional[LayoutManifest]:
        """
        Analyze a single slide layout.

        Args:
            layout: python-pptx SlideLayout object
            index: Layout index
            template_name: Name of the source template

        Returns:
            LayoutManifest or None if analysis fails
        """
        try:
            layout_id = f"LAYOUT_{index:02d}"
            layout_name = layout.name or f"Layout {index}"

            # Extract placeholders
            placeholders = []
            body_count = 0
            title_count = 0
            picture_count = 0

            for shape in layout.shapes:
                if shape.is_placeholder:
                    placeholder_info = self._extract_placeholder_info(shape)
                    if placeholder_info:
                        placeholders.append(placeholder_info)

                        if placeholder_info["type"] == "title":
                            title_count += 1
                        elif placeholder_info["type"] == "body":
                            body_count += 1
                        elif placeholder_info["type"] == "picture":
                            picture_count += 1

            # Estimate element count range based on placeholders
            element_count = max(body_count, 1)
            element_count_range = (max(1, element_count - 1), element_count + 2)

            # Infer category from layout name and structure
            category = self._infer_category(layout_name, element_count, placeholders)

            # Generate basic description
            description = self._generate_basic_description(
                layout_name, category, element_count, placeholders
            )

            # Extract keywords from layout name
            keywords = self._extract_keywords(layout_name)

            return LayoutManifest(
                layout_id=layout_id,
                layout_name=layout_name,
                source_template=template_name,
                description=description,
                category=category,
                element_count_range=element_count_range,
                placeholders=placeholders,
                keywords=keywords,
            )

        except Exception as e:
            logger.error(f"Failed to analyze layout {index}: {e}")
            return None

    def _extract_placeholder_info(self, shape) -> Dict[str, Any]:
        """Extract placeholder information from a shape"""
        try:
            # Get placeholder type
            ph_type = "body"
            if hasattr(shape, 'placeholder_format'):
                ph_format = shape.placeholder_format
                if ph_format.type == 1:  # Title
                    ph_type = "title"
                elif ph_format.type == 2:  # Body
                    ph_type = "body"
                elif ph_format.type == 18:  # Picture
                    ph_type = "picture"
                elif ph_format.type == 6:  # Chart
                    ph_type = "chart"
                elif ph_format.type == 7:  # Table
                    ph_type = "table"

            # Get position
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            # Convert EMUs to relative positions (0-1)
            # Note: These would need slide dimensions for proper conversion
            position = [left, top, width, height]

            return {
                "id": shape.shape_id if hasattr(shape, 'shape_id') else 0,
                "type": ph_type,
                "position": position,
                "name": shape.name or "",
            }

        except Exception as e:
            logger.debug(f"Failed to extract placeholder info: {e}")
            return None

    def _infer_category(
        self,
        layout_name: str,
        element_count: int,
        placeholders: List[Dict]
    ) -> str:
        """Infer layout category from name and structure"""
        name_lower = layout_name.lower()

        # === Step 1: Keyword matching (high priority) ===

        # Matrix patterns
        if any(kw in name_lower for kw in ["matrix", "quad", "swot", "四象限", "矩阵", "2x2", "3x3", "九宫格"]):
            return "MATRIX"

        # Process patterns
        if any(kw in name_lower for kw in ["process", "flow", "step", "流程", "步骤", "顺序", "阶段", "环节"]):
            return "PROCESS"

        # Timeline patterns
        if any(kw in name_lower for kw in ["timeline", "gantt", "schedule", "时间", "进度", "甘特", "里程碑", "milestone"]):
            return "TIMELINE"

        # Table/comparison patterns
        if any(kw in name_lower for kw in ["table", "comparison", "对比", "表格", "比较", "优劣", "方案"]):
            return "TABLE"

        # Data viz patterns
        if any(kw in name_lower for kw in ["chart", "data", "graph", "图表", "数据", "统计", "指标", "radar", "pie", "bar"]):
            return "DATA_VIZ"

        # Key insight patterns
        if element_count == 1 or any(kw in name_lower for kw in ["insight", "key", "highlight", "核心", "重点", "关键", "强调", "总结", "summary"]):
            return "KEY_INSIGHT"

        # === Step 2: Structure analysis ===

        # Analyze placeholder positions to infer layout type
        if placeholders and len(placeholders) >= 2:
            # Get positions (normalized to 0-1 range)
            body_placeholders = [p for p in placeholders if p.get("type") == "body"]

            if len(body_placeholders) >= 2:
                positions = []
                for p in body_placeholders:
                    pos = p.get("position", [0, 0, 0, 0])
                    # Normalize position (EMUs to relative, assume 9144000 EMUs = 1 inch, slide is ~10x7.5 inches)
                    left = pos[0] / 9144000 / 10  # normalize to 0-1
                    top = pos[1] / 9144000 / 7.5
                    width = pos[2] / 9144000 / 10
                    height = pos[3] / 9144000 / 7.5
                    positions.append({"left": left, "top": top, "width": width, "height": height})

                # Check for matrix (2x2 grid pattern)
                if len(positions) == 4:
                    # Sort by top-left position
                    sorted_pos = sorted(positions, key=lambda p: (p["top"], p["left"]))
                    # Check if roughly in 2x2 grid
                    top_row = [p for p in sorted_pos if p["top"] < 0.5]
                    bottom_row = [p for p in sorted_pos if p["top"] >= 0.5]
                    if len(top_row) == 2 and len(bottom_row) == 2:
                        return "MATRIX"

                # Check for process (horizontal sequence)
                if len(positions) >= 3:
                    # Check if placeholders are arranged horizontally
                    left_values = [p["left"] for p in positions]
                    top_values = [p["top"] for p in positions]

                    # If horizontal spread > vertical spread, likely process
                    h_spread = max(left_values) - min(left_values)
                    v_spread = max(top_values) - min(top_values)

                    if h_spread > 0.4 and h_spread > v_spread * 1.5:
                        return "PROCESS"

                    # If vertical spread > horizontal spread, could be timeline or process
                    if v_spread > 0.4 and v_spread > h_spread * 1.5:
                        return "TIMELINE"

                # Check for timeline (elements along a line with time indicators)
                if len(positions) >= 3:
                    # Timeline often has elements along bottom or center
                    avg_top = sum(p["top"] for p in positions) / len(positions)
                    if avg_top > 0.5:  # Elements in lower half
                        return "TIMELINE"

        # === Step 3: Element count heuristics ===

        # Single element = key insight
        if element_count == 1:
            return "KEY_INSIGHT"

        # 2 elements = comparison/table
        if element_count == 2:
            return "TABLE"

        # 4 elements could be matrix or parallel
        if element_count == 4:
            # Default to parallel unless structure suggests matrix
            return "PARALLEL"

        # 9 elements = 3x3 matrix
        if element_count == 9:
            return "MATRIX"

        # Default to parallel
        return "PARALLEL"

    def _generate_basic_description(
        self,
        layout_name: str,
        category: str,
        element_count: int,
        placeholders: List[Dict]
    ) -> str:
        """Generate a basic description of the layout"""
        category_descriptions = {
            "MATRIX": f"矩阵布局，适合四象限或多维度分析",
            "PROCESS": f"流程布局，适合展示{element_count}步顺序流程",
            "TIMELINE": f"时间线布局，适合展示项目进度和里程碑",
            "TABLE": f"表格布局，适合对比分析",
            "DATA_VIZ": f"数据可视化布局，适合展示图表和数据",
            "KEY_INSIGHT": f"核心观点布局，适合强调单一关键信息",
            "PARALLEL": f"并列布局，适合{element_count}个并列要点展示",
        }

        base_desc = category_descriptions.get(category, f"适合{element_count}要素内容展示")

        # Add placeholder info
        title_count = sum(1 for p in placeholders if p.get("type") == "title")
        body_count = sum(1 for p in placeholders if p.get("type") == "body")
        picture_count = sum(1 for p in placeholders if p.get("type") == "picture")

        extras = []
        if picture_count > 0:
            extras.append(f"{picture_count}个图片占位符")
        if title_count > 1:
            extras.append(f"{title_count}个标题区")

        if extras:
            base_desc += f"，包含{', '.join(extras)}"

        return base_desc

    def _extract_keywords(self, layout_name: str) -> List[str]:
        """Extract keywords from layout name"""
        keywords = []

        # Common keywords to look for
        keyword_map = {
            "流程": ["流程", "步骤", "阶段"],
            "对比": ["对比", "比较", "优劣"],
            "时间": ["时间", "进度", "计划"],
            "矩阵": ["矩阵", "象限", "SWOT"],
            "数据": ["数据", "图表", "统计"],
            "核心": ["核心", "关键", "重点"],
        }

        name_lower = layout_name.lower()
        for key, kws in keyword_map.items():
            for kw in kws:
                if kw in name_lower:
                    keywords.append(key)
                    break

        return list(set(keywords))

    def _generate_ai_descriptions(self, manifests: List[LayoutManifest]) -> List[LayoutManifest]:
        """Use AI to generate semantic descriptions for layouts"""
        try:
            # Import AI service if available
            from lib.report_workflow.ai_service import report_ai_service

            if not report_ai_service.is_configured():
                logger.warning("AI service not configured, using basic descriptions")
                return manifests

            # For each manifest, generate enhanced description
            for manifest in manifests:
                prompt = f"""分析以下PPT布局并生成语义描述：

布局名称：{manifest.layout_name}
类别：{manifest.category}
要素数量：{manifest.element_count_range[0]}-{manifest.element_count_range[1]}
占位符数量：{len(manifest.placeholders)}

请用一句话描述这个布局最适合什么场景使用（20字以内）。

只输出描述文字，不要其他内容。"""

                try:
                    import asyncio
                    description = asyncio.run(
                        report_ai_service._call_api("你是一位PPT设计专家", prompt)
                    )
                    if description:
                        manifest.description = description.strip()[:100]
                except Exception as e:
                    logger.debug(f"AI description generation failed: {e}")

            return manifests

        except ImportError:
            logger.warning("AI service not available, using basic descriptions")
            return manifests
        except Exception as e:
            logger.warning(f"AI description generation failed: {e}")
            return manifests

    def save_manifests(
        self,
        manifests: List[LayoutManifest],
        output_path: str
    ) -> bool:
        """
        Save layout manifests to a JSON file.

        Args:
            manifests: List of layout manifests
            output_path: Path to save the JSON file

        Returns:
            True if successful, False otherwise
        """
        try:
            data = {
                "version": "1.0",
                "layouts": [m.to_dict() for m in manifests]
            }

            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)

            logger.info(f"Saved {len(manifests)} layout manifests to {output_path}")
            return True

        except Exception as e:
            logger.error(f"Failed to save manifests: {e}")
            return False


# Singleton instance
_analyzer = None


def get_template_analyzer() -> TemplateAnalyzer:
    """Get the singleton template analyzer instance"""
    global _analyzer
    if _analyzer is None:
        _analyzer = TemplateAnalyzer()
    return _analyzer

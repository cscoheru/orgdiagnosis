"""
Visual Elements Factory for PPTX Rendering

Provides shape generators, chart builders, and visual element utilities
for creating professional PowerPoint slides.

Created: 2026-03-23
"""

from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from loguru import logger

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import math


@dataclass
class Position:
    """Shape position in inches"""
    left: float
    top: float
    width: float
    height: float

    def to_emu(self) -> Tuple[Emu, Emu, Emu, Emu]:
        """Convert to EMU units for python-pptx"""
        return (
            Inches(self.left),
            Inches(self.top),
            Inches(self.width),
            Inches(self.height)
        )


@dataclass
class ColorScheme:
    """Color scheme for a slide or theme"""
    primary: RGBColor
    secondary: RGBColor
    accent: RGBColor
    text_dark: RGBColor
    text_light: RGBColor
    background: RGBColor
    white: RGBColor = RGBColor(255, 255, 255)

    @classmethod
    def from_dict(cls, data: Dict[str, str]) -> "ColorScheme":
        """Create from hex color dict"""
        def hex_to_rgb(hex_str: str) -> RGBColor:
            hex_str = hex_str.lstrip('#')
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

        return cls(
            primary=hex_to_rgb(data.get("primary", "#00528B")),
            secondary=hex_to_rgb(data.get("secondary", "#336699")),
            accent=hex_to_rgb(data.get("accent", "#E74C3C")),
            text_dark=hex_to_rgb(data.get("text_dark", "#333333")),
            text_light=hex_to_rgb(data.get("text_light", "#808080")),
            background=hex_to_rgb(data.get("background", "#F5F5F5"))
        )


class ShapeFactory:
    """
    Factory for creating PowerPoint shapes.

    Provides methods to create various shapes with consistent styling.
    """

    @staticmethod
    def add_rounded_rectangle(
        slide,
        position: Position,
        fill_color: RGBColor,
        text: str = "",
        font_size: int = 14,
        font_color: RGBColor = None,
        corner_radius: float = 0.1
    ) -> Any:
        """Add a rounded rectangle shape"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            *position.to_emu()
        )

        # Fill
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color

        # No line
        shape.line.fill.background()

        # Text
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color or RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].space_before = Pt(0)
            tf.paragraphs[0].space_after = Pt(0)

        return shape

    @staticmethod
    def add_card(
        slide,
        position: Position,
        title: str,
        content: str,
        colors: ColorScheme,
        icon: str = None,
        number: str = None
    ) -> Any:
        """Add a card-style element with title and content"""
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            *position.to_emu()
        )
        card.fill.solid()
        card.fill.fore_color.rgb = colors.white
        card.line.color.rgb = colors.background
        card.line.width = Pt(1)

        # Shadow effect (simulated with offset)
        shadow = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(position.left + 0.05),
            Inches(position.top + 0.05),
            Inches(position.width),
            Inches(position.height)
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(220, 220, 220)
        shadow.line.fill.background()

        # Move shadow behind card
        spTree = slide.shapes._spTree
        sp = shadow._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # Number badge (optional)
        if number:
            badge_pos = Position(position.left + 0.1, position.top + 0.1, 0.4, 0.4)
            ShapeFactory.add_circle(slide, badge_pos, colors.primary, number, font_size=12)

        # Title
        title_top = position.top + (0.6 if number else 0.2)
        title_box = slide.shapes.add_textbox(
            Inches(position.left + 0.2),
            Inches(title_top),
            Inches(position.width - 0.4),
            Inches(0.4)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = colors.text_dark

        # Content
        content_top = title_top + 0.5
        content_box = slide.shapes.add_textbox(
            Inches(position.left + 0.2),
            Inches(content_top),
            Inches(position.width - 0.4),
            Inches(position.height - (content_top - position.top) - 0.2)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(11)
        p.font.color.rgb = colors.text_light

        return card

    @staticmethod
    def add_circle(
        slide,
        position: Position,
        fill_color: RGBColor,
        text: str = "",
        font_size: int = 12,
        font_color: RGBColor = None
    ) -> Any:
        """Add a circle shape"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            *position.to_emu()
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()

        if text:
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color or RGBColor(255, 255, 255)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        return shape

    @staticmethod
    def add_arrow(
        slide,
        start: Tuple[float, float],
        end: Tuple[float, float],
        color: RGBColor,
        thickness: float = 2.0
    ) -> Any:
        """Add an arrow connector"""
        # Calculate position
        left = min(start[0], end[0])
        top = min(start[1], end[1])
        width = abs(end[0] - start[0])
        height = abs(end[1] - start[1])

        if width < 0.1:
            width = 0.1
        if height < 0.1:
            height = 0.1

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW if width > height else MSO_SHAPE.DOWN_ARROW,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(thickness / 72)  # Convert pt to inches
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        return shape

    @staticmethod
    def add_hexagon(
        slide,
        position: Position,
        fill_color: RGBColor,
        text: str = "",
        font_size: int = 12
    ) -> Any:
        """Add a hexagon shape"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.HEXAGON,
            *position.to_emu()
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(2)

        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

        return shape


class ChartBuilder:
    """
    Builder for creating chart-like visualizations.

    Since python-pptx has limited chart support, we build charts
    using shapes for more flexibility.
    """

    @staticmethod
    def build_matrix_2x2(
        slide,
        position: Position,
        quadrants: List[Dict[str, str]],
        colors: ColorScheme,
        x_label: str = "",
        y_label: str = ""
    ) -> List[Any]:
        """
        Build a 2x2 matrix (like SWOT).

        Args:
            quadrants: List of 4 dicts with 'title' and 'content' keys
                      Order: [top-left, top-right, bottom-left, bottom-right]
        """
        shapes = []
        cell_width = position.width / 2 - 0.1
        cell_height = position.height / 2 - 0.1

        # Quadrant colors
        quad_colors = [
            colors.primary,
            colors.secondary,
            colors.accent,
            RGBColor(100, 150, 100)  # Green-ish
        ]

        positions = [
            (position.left, position.top),  # top-left
            (position.left + position.width / 2 + 0.1, position.top),  # top-right
            (position.left, position.top + position.height / 2 + 0.1),  # bottom-left
            (position.left + position.width / 2 + 0.1, position.top + position.height / 2 + 0.1),  # bottom-right
        ]

        for i, (q_pos, quad, color) in enumerate(zip(positions, quadrants[:4], quad_colors)):
            cell_pos = Position(q_pos[0], q_pos[1], cell_width, cell_height)

            # Background
            bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                *cell_pos.to_emu()
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = color
            bg.line.fill.background()
            shapes.append(bg)

            # Title
            title_box = slide.shapes.add_textbox(
                Inches(q_pos[0] + 0.2),
                Inches(q_pos[1] + 0.15),
                Inches(cell_width - 0.4),
                Inches(0.4)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = quad.get("title", "")
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            shapes.append(title_box)

            # Content
            content_box = slide.shapes.add_textbox(
                Inches(q_pos[0] + 0.2),
                Inches(q_pos[1] + 0.6),
                Inches(cell_width - 0.4),
                Inches(cell_height - 0.8)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = quad.get("content", "")
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            shapes.append(content_box)

        # Axis labels
        if x_label:
            label_box = slide.shapes.add_textbox(
                Inches(position.left),
                Inches(position.top + position.height + 0.1),
                Inches(position.width),
                Inches(0.3)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = x_label
            p.font.size = Pt(10)
            p.font.color.rgb = colors.text_light
            p.alignment = PP_ALIGN.CENTER
            shapes.append(label_box)

        return shapes

    @staticmethod
    def build_process_flow(
        slide,
        position: Position,
        steps: List[Dict[str, str]],
        colors: ColorScheme,
        orientation: str = "horizontal"
    ) -> List[Any]:
        """
        Build a process flow diagram.

        Args:
            steps: List of dicts with 'title' and 'description' keys
            orientation: 'horizontal' or 'vertical'
        """
        shapes = []
        n_steps = len(steps)

        if orientation == "horizontal":
            step_width = (position.width - (n_steps - 1) * 0.3) / n_steps
            step_height = min(1.2, position.height)

            for i, step in enumerate(steps):
                x = position.left + i * (step_width + 0.3)
                y = position.top + (position.height - step_height) / 2

                # Step box
                box_pos = Position(x, y, step_width, step_height)
                box = ShapeFactory.add_rounded_rectangle(
                    slide, box_pos, colors.primary,
                    text="", font_size=12
                )
                shapes.append(box)

                # Step number
                num_pos = Position(x + step_width / 2 - 0.25, y - 0.15, 0.5, 0.5)
                ShapeFactory.add_circle(slide, num_pos, colors.accent, str(i + 1), font_size=14)

                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(x + 0.1),
                    Inches(y + 0.3),
                    Inches(step_width - 0.2),
                    Inches(0.4)
                )
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = step.get("title", "")
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
                shapes.append(title_box)

                # Arrow between steps
                if i < n_steps - 1:
                    arrow_x = x + step_width
                    arrow_y = y + step_height / 2
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        Inches(arrow_x),
                        Inches(arrow_y - 0.1),
                        Inches(0.25),
                        Inches(0.2)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = colors.accent
                    arrow.line.fill.background()
                    shapes.append(arrow)

        else:  # vertical
            step_width = position.width
            step_height = (position.height - (n_steps - 1) * 0.2) / n_steps

            for i, step in enumerate(steps):
                x = position.left
                y = position.top + i * (step_height + 0.2)

                box_pos = Position(x, y, step_width, step_height)
                box = ShapeFactory.add_rounded_rectangle(
                    slide, box_pos, colors.primary,
                    text=step.get("title", ""), font_size=14
                )
                shapes.append(box)

                # Arrow
                if i < n_steps - 1:
                    arrow_y = y + step_height
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW,
                        Inches(x + step_width / 2 - 0.15),
                        Inches(arrow_y),
                        Inches(0.3),
                        Inches(0.15)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = colors.accent
                    arrow.line.fill.background()
                    shapes.append(arrow)

        return shapes

    @staticmethod
    def build_timeline(
        slide,
        position: Position,
        milestones: List[Dict[str, str]],
        colors: ColorScheme
    ) -> List[Any]:
        """
        Build a timeline visualization.

        Args:
            milestones: List of dicts with 'date', 'title', 'description' keys
        """
        shapes = []
        n_items = len(milestones)

        if n_items == 0:
            return shapes

        # Timeline line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(position.left),
            Inches(position.top + 0.4),
            Inches(position.width),
            Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = colors.secondary
        line.line.fill.background()
        shapes.append(line)

        # Milestones
        spacing = position.width / (n_items + 1)

        for i, ms in enumerate(milestones):
            x = position.left + (i + 1) * spacing

            # Milestone dot
            dot_pos = Position(x - 0.15, position.top + 0.25, 0.3, 0.3)
            ShapeFactory.add_circle(slide, dot_pos, colors.primary, str(i + 1), font_size=10)

            # Date label (above line)
            date_box = slide.shapes.add_textbox(
                Inches(x - 0.5),
                Inches(position.top - 0.1),
                Inches(1.0),
                Inches(0.3)
            )
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = ms.get("date", "")
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = colors.text_dark
            p.alignment = PP_ALIGN.CENTER
            shapes.append(date_box)

            # Title (below line)
            title_box = slide.shapes.add_textbox(
                Inches(x - 0.8),
                Inches(position.top + 0.7),
                Inches(1.6),
                Inches(0.8)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = ms.get("title", "")
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = colors.text_dark
            p.alignment = PP_ALIGN.CENTER
            shapes.append(title_box)

        return shapes

    @staticmethod
    def build_radar_chart(
        slide,
        position: Position,
        dimensions: List[Dict[str, Any]],
        colors: ColorScheme
    ) -> List[Any]:
        """
        Build a radar/spider chart using shapes.

        Args:
            dimensions: List of dicts with 'name' and 'value' (0-100) keys
        """
        shapes = []
        n_dims = len(dimensions)

        if n_dims < 3:
            return shapes

        # Center and radius
        center_x = position.left + position.width / 2
        center_y = position.top + position.height / 2
        radius = min(position.width, position.height) / 2 - 0.5

        # Draw axis lines and labels
        for i, dim in enumerate(dimensions):
            angle = (2 * math.pi * i / n_dims) - math.pi / 2  # Start from top

            # End point
            end_x = center_x + radius * math.cos(angle)
            end_y = center_y + radius * math.sin(angle)

            # Axis line
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(center_x),
                Inches(center_y),
                Inches(end_x),
                Inches(end_y)
            )
            line.line.color.rgb = colors.text_light
            line.line.width = Pt(1)
            shapes.append(line)

            # Label
            label_x = center_x + (radius + 0.3) * math.cos(angle) - 0.4
            label_y = center_y + (radius + 0.3) * math.sin(angle) - 0.15

            label_box = slide.shapes.add_textbox(
                Inches(label_x),
                Inches(label_y),
                Inches(0.8),
                Inches(0.3)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = dim.get("name", "")
            p.font.size = Pt(10)
            p.font.color.rgb = colors.text_dark
            p.alignment = PP_ALIGN.CENTER
            shapes.append(label_box)

        # Draw data polygon (simplified as a series of connected lines)
        # For a full implementation, you'd use a freeform shape
        # Here we just show the data points
        for i, dim in enumerate(dimensions):
            value = dim.get("value", 50) / 100  # Normalize to 0-1
            angle = (2 * math.pi * i / n_dims) - math.pi / 2

            point_x = center_x + radius * value * math.cos(angle)
            point_y = center_y + radius * value * math.sin(angle)

            # Data point
            dot_pos = Position(point_x - 0.1, point_y - 0.1, 0.2, 0.2)
            ShapeFactory.add_circle(slide, dot_pos, colors.accent, "", font_size=8)

        return shapes

    @staticmethod
    def build_pyramid(
        slide,
        position: Position,
        levels: List[Dict[str, str]],
        colors: ColorScheme
    ) -> List[Any]:
        """
        Build a pyramid hierarchy diagram.

        Args:
            levels: List of dicts from top to bottom with 'title' and 'description'
        """
        shapes = []
        n_levels = len(levels)

        if n_levels == 0:
            return shapes

        level_height = position.height / n_levels

        for i, level in enumerate(levels):
            # Width increases towards bottom
            width_ratio = (i + 1) / n_levels
            level_width = position.width * (0.4 + 0.6 * width_ratio)
            x_offset = (position.width - level_width) / 2

            y = position.top + i * level_height
            x = position.left + x_offset

            # Trapezoid approximation using chevron or rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON if i == 0 else MSO_SHAPE.PENTAGON,
                Inches(x),
                Inches(y),
                Inches(level_width),
                Inches(level_height - 0.05)
            )
            shape.fill.solid()
            # Gradient color from primary to secondary
            color_ratio = i / (n_levels - 1) if n_levels > 1 else 0
            r = int(colors.primary.red * (1 - color_ratio) + colors.secondary.red * color_ratio)
            g = int(colors.primary.green * (1 - color_ratio) + colors.secondary.green * color_ratio)
            b = int(colors.primary.blue * (1 - color_ratio) + colors.secondary.blue * color_ratio)
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.fill.background()
            shapes.append(shape)

            # Text
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = level.get("title", "")
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        return shapes

    @staticmethod
    def build_comparison_table(
        slide,
        position: Position,
        headers: List[str],
        rows: List[List[str]],
        colors: ColorScheme
    ) -> List[Any]:
        """
        Build a comparison table using shapes.

        Args:
            headers: Column headers
            rows: Table data as list of rows
        """
        shapes = []
        n_cols = len(headers)
        n_rows = len(rows) + 1  # +1 for header

        if n_cols == 0 or n_rows == 1:
            return shapes

        col_width = position.width / n_cols
        row_height = position.height / n_rows

        # Header row
        for j, header in enumerate(headers):
            x = position.left + j * col_width
            y = position.top

            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(col_width - 0.02),
                Inches(row_height - 0.02)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors.primary
            cell.line.color.rgb = colors.white
            cell.line.width = Pt(1)
            shapes.append(cell)

            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.text = header
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        # Data rows
        for i, row in enumerate(rows):
            for j, cell_text in enumerate(row):
                x = position.left + j * col_width
                y = position.top + (i + 1) * row_height

                bg_color = colors.white if i % 2 == 0 else colors.background
                cell = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x),
                    Inches(y),
                    Inches(col_width - 0.02),
                    Inches(row_height - 0.02)
                )
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                cell.line.color.rgb = colors.background
                cell.line.width = Pt(1)
                shapes.append(cell)

                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = str(cell_text)
                p.font.size = Pt(11)
                p.font.color.rgb = colors.text_dark
                p.alignment = PP_ALIGN.CENTER

        return shapes


class LayoutRenderer:
    """
    Renders layouts based on layout configuration.

    Maps layout IDs to actual rendering methods.
    """

    def __init__(self, colors: ColorScheme = None):
        """Initialize with color scheme"""
        self.colors = colors or ColorScheme(
            primary=RGBColor(0, 82, 139),
            secondary=RGBColor(51, 102, 153),
            accent=RGBColor(231, 76, 60),
            text_dark=RGBColor(51, 51, 51),
            text_light=RGBColor(128, 128, 128),
            background=RGBColor(245, 245, 245)
        )

    def render(
        self,
        slide,
        layout_id: str,
        content: Dict[str, Any],
        position: Position = None
    ) -> List[Any]:
        """
        Render a layout on the slide.

        Args:
            slide: python-pptx slide object
            layout_id: Layout identifier (e.g., "PARALLEL_03_CARDS")
            content: Content dict with bullets, title, etc.
            position: Rendering position (default: full slide)
        """
        if position is None:
            position = Position(0.5, 1.5, 12.333, 5.5)  # Default content area

        # Parse layout ID
        parts = layout_id.split("_")
        category = parts[0] if parts else "PARALLEL"

        # Route to appropriate renderer
        if category == "KEY" or layout_id == "KEY_INSIGHT_01":
            return self._render_key_insight(slide, content, position)
        elif category == "PARALLEL" or category == "BULLET":
            return self._render_parallel(slide, layout_id, content, position)
        elif category == "MATRIX":
            return self._render_matrix(slide, layout_id, content, position)
        elif category == "PROCESS":
            return self._render_process(slide, layout_id, content, position)
        elif category == "TIMELINE":
            return self._render_timeline(slide, content, position)
        elif category == "TABLE":
            return self._render_table(slide, content, position)
        elif category == "HIERARCHY":
            return self._render_hierarchy(slide, content, position)
        else:
            # Fallback to bullet points
            return self._render_bullet_list(slide, content, position)

    def _render_key_insight(self, slide, content: Dict, position: Position) -> List[Any]:
        """Render a key insight slide"""
        shapes = []

        # Large central text
        key_message = content.get("key_message", "")
        title = content.get("title", "")

        # Title at top
        title_box = slide.shapes.add_textbox(
            Inches(position.left),
            Inches(0.5),
            Inches(position.width),
            Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors.text_dark
        p.alignment = PP_ALIGN.CENTER
        shapes.append(title_box)

        # Key message - large centered
        msg_box = slide.shapes.add_textbox(
            Inches(position.left + 0.5),
            Inches(2.5),
            Inches(position.width - 1),
            Inches(2.0)
        )
        tf = msg_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = key_message
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.colors.primary
        p.alignment = PP_ALIGN.CENTER
        shapes.append(msg_box)

        return shapes

    def _render_parallel(self, slide, layout_id: str, content: Dict, position: Position) -> List[Any]:
        """Render parallel/cards layout"""
        shapes = []
        bullets = content.get("bullets", [])
        n_items = len(bullets)

        if n_items == 0:
            return shapes

        # Determine card layout based on count
        if n_items <= 3:
            # Horizontal cards
            card_width = (position.width - (n_items - 1) * 0.2) / n_items
            card_height = position.height

            for i, bullet in enumerate(bullets):
                x = position.left + i * (card_width + 0.2)
                card_pos = Position(x, position.top, card_width, card_height)

                # Parse bullet for title/content
                parts = bullet.split("：", 1) if "：" in bullet else bullet.split(":", 1)
                if len(parts) == 2:
                    title, desc = parts
                else:
                    title = bullet[:30]
                    desc = ""

                ShapeFactory.add_card(slide, card_pos, title, desc, self.colors, number=str(i + 1))

        elif n_items <= 6:
            # 2-row grid
            n_cols = 3
            n_rows = (n_items + 2) // 3
            card_width = (position.width - 2 * 0.15) / 3
            card_height = (position.height - (n_rows - 1) * 0.15) / n_rows

            for i, bullet in enumerate(bullets):
                row = i // 3
                col = i % 3
                x = position.left + col * (card_width + 0.15)
                y = position.top + row * (card_height + 0.15)
                card_pos = Position(x, y, card_width, card_height)

                parts = bullet.split("：", 1) if "：" in bullet else bullet.split(":", 1)
                if len(parts) == 2:
                    title, desc = parts
                else:
                    title = bullet[:30]
                    desc = ""

                ShapeFactory.add_card(slide, card_pos, title, desc, self.colors, number=str(i + 1))

        else:
            # More than 6 - use compact list
            self._render_bullet_list(slide, content, position)

        return shapes

    def _render_matrix(self, slide, layout_id: str, content: Dict, position: Position) -> List[Any]:
        """Render matrix layout"""
        bullets = content.get("bullets", [])

        # Try to parse as 2x2 or 3x3
        if len(bullets) == 4:
            quadrants = [{"title": b[:20], "content": b} for b in bullets]
            return ChartBuilder.build_matrix_2x2(slide, position, quadrants, self.colors)
        else:
            # Fallback to grid of cards
            return self._render_parallel(slide, "PARALLEL_04", content, position)

    def _render_process(self, slide, layout_id: str, content: Dict, position: Position) -> List[Any]:
        """Render process flow layout"""
        bullets = content.get("bullets", [])
        steps = [{"title": b} for b in bullets]

        # Determine orientation from layout ID
        orientation = "vertical" if "V" in layout_id else "horizontal"

        return ChartBuilder.build_process_flow(slide, position, steps, self.colors, orientation)

    def _render_timeline(self, slide, content: Dict, position: Position) -> List[Any]:
        """Render timeline layout"""
        bullets = content.get("bullets", [])
        milestones = [{"title": b, "date": f"阶段{i+1}"} for i, b in enumerate(bullets)]

        return ChartBuilder.build_timeline(slide, position, milestones, self.colors)

    def _render_table(self, slide, content: Dict, position: Position) -> List[Any]:
        """Render comparison table"""
        bullets = content.get("bullets", [])

        # Simple 2-column comparison
        headers = ["维度", "描述"]
        rows = []
        for b in bullets:
            parts = b.split("：", 1) if "：" in b else b.split(":", 1)
            if len(parts) == 2:
                rows.append([parts[0], parts[1]])
            else:
                rows.append(["", b])

        return ChartBuilder.build_comparison_table(slide, position, headers, rows, self.colors)

    def _render_hierarchy(self, slide, content: Dict, position: Position) -> List[Any]:
        """Render hierarchy/pyramid layout"""
        bullets = content.get("bullets", [])
        levels = [{"title": b} for b in bullets]

        return ChartBuilder.build_pyramid(slide, position, levels, self.colors)

    def _render_bullet_list(self, slide, content: Dict, position: Position) -> List[Any]:
        """Render simple bullet list"""
        shapes = []
        bullets = content.get("bullets", [])

        bullet_text = "\n\n".join([f"• {b}" for b in bullets])

        text_box = slide.shapes.add_textbox(
            Inches(position.left),
            Inches(position.top),
            Inches(position.width),
            Inches(position.height)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet_text
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors.text_dark
        shapes.append(text_box)

        return shapes


# Convenience functions
def get_layout_renderer(theme_colors: Dict[str, str] = None) -> LayoutRenderer:
    """Get a layout renderer with optional theme colors"""
    if theme_colors:
        colors = ColorScheme.from_dict(theme_colors)
        return LayoutRenderer(colors)
    return LayoutRenderer()

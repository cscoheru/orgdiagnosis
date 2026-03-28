"""
Competency Co-pilot — 源材料转换脚本

将 docx/xlsx/pptx/pdf 文件批量转换为 source_materials.json

用法:
  cd backend
  python scripts/competency/convert_files.py /path/to/your/files

  # 或指定输出路径
  python scripts/competency/convert_files.py /path/to/your/files -o scripts/competency/source_materials.json

  # 只处理特定类型
  python scripts/competency/convert_files.py /path/to/your/files --types docx pdf
"""

import argparse
import json
import sys
from pathlib import Path
from typing import List, Dict

SUPPORTED_TYPES = {'.docx', '.xlsx', '.pptx', '.pdf', '.txt', '.md'}


def extract_docx(filepath: Path) -> str:
    """提取 Word 文档文本"""
    from docx import Document
    doc = Document(str(filepath))

    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # 也提取表格内容
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    return '\n'.join(parts)


def extract_xlsx(filepath: Path) -> str:
    """提取 Excel 表格文本"""
    from openpyxl import load_workbook
    wb = load_workbook(str(filepath), read_only=True, data_only=True)

    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"[工作表: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            row_text = ' | '.join(str(c).strip() if c else '' for c in row)
            row_text = row_text.strip(' |')
            if row_text:
                parts.append(row_text)

    return '\n'.join(parts)


def extract_pptx(filepath: Path) -> str:
    """提取 PPT 幻灯片文本"""
    from pptx import Presentation
    prs = Presentation(str(filepath))

    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            # 也提取表格
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_texts.append(row_text)
        if slide_texts:
            parts.append(f"[幻灯片 {i}]")
            parts.extend(slide_texts)

    return '\n'.join(parts)


def extract_pdf(filepath: Path) -> str:
    """提取 PDF 文本"""
    import fitz  # pymupdf
    doc = fitz.open(str(filepath))

    parts = []
    for page in doc:
        text = page.get_text().strip()
        if text:
            parts.append(text)

    doc.close()
    return '\n'.join(parts)


def extract_txt(filepath: Path) -> str:
    """读取纯文本文件"""
    encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16']
    for enc in encodings:
        try:
            return filepath.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return filepath.read_text(encoding='utf-8', errors='replace')


def extract_file(filepath: Path) -> Dict:
    """根据文件类型提取文本"""
    ext = filepath.suffix.lower()

    extractors = {
        '.docx': extract_docx,
        '.xlsx': extract_xlsx,
        '.pptx': extract_pptx,
        '.pdf': extract_pdf,
        '.txt': extract_txt,
        '.md': extract_txt,
    }

    extractor = extractors.get(ext)
    if not extractor:
        return None

    content = extractor(filepath)
    return {
        "source": filepath.name,
        "content": content,
    }


def main():
    parser = argparse.ArgumentParser(description="将文件转换为 source_materials.json")
    parser.add_argument("input_dir", type=Path, help="包含源文件的目录路径")
    parser.add_argument(
        "-o", "--output",
        type=Path,
        default=Path(__file__).parent / "source_materials.json",
        help="输出 JSON 文件路径",
    )
    parser.add_argument(
        "--types",
        nargs='+',
        default=None,
        help="只处理指定类型，如: docx pdf",
    )
    args = parser.parse_args()

    input_dir = args.input_dir
    if not input_dir.exists():
        print(f"错误: 目录不存在: {input_dir}")
        sys.exit(1)

    # 收集文件
    allowed_types = {f'.{t}' for t in args.types} if args.types else SUPPORTED_TYPES
    files = sorted([
        f for f in input_dir.iterdir()
        if f.is_file() and f.suffix.lower() in allowed_types
    ])

    if not files:
        print(f"在 {input_dir} 中未找到支持的文件类型: {', '.join(allowed_types)}")
        sys.exit(1)

    print(f"找到 {len(files)} 个文件待处理:")
    for f in files:
        print(f"  - {f.name} ({f.suffix})")

    # 提取
    results = []
    for f in files:
        print(f"\n处理: {f.name}...", end=' ')
        try:
            item = extract_file(f)
            if item and item['content'].strip():
                # 截断过长内容 (避免 AI prompt 超限)
                if len(item['content']) > 15000:
                    item['content'] = item['content'][:15000] + '\n\n[... 内容已截断 ...]'
                results.append(item)
                print(f"✓ ({len(item['content'])} 字符)")
            else:
                print("⚠ 跳过 (无内容)")
        except Exception as e:
            print(f"✗ 错误: {e}")

    # 保存
    output_path = args.output
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)

    print(f"\n{'='*50}")
    print(f"已转换 {len(results)} 个文件")
    print(f"输出: {output_path}")
    print(f"总字符数: {sum(len(r['content']) for r in results)}")


if __name__ == "__main__":
    main()
